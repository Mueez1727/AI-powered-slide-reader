"""
PPT/PPTX extraction service using python-pptx.
Extracts structured slide data — heading, main content, table text,
speaker notes, and OCR text from embedded images.
"""

from __future__ import annotations

import logging
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.image_captioning_service import generate_image_caption
from app.services.ocr_service import MAX_IMAGES_PER_SLIDE, ocr_image_bytes
from app.services.text_cleaner import clean_text

logger = logging.getLogger(__name__)


class PPTExtractor:
    """Extract structured slide data from PowerPoint files."""

    # ── public API ──────────────────────────────────────────

    @staticmethod
    def extract(file_path: str) -> list[dict]:
        """
        Extract every slide of a PPTX file.

        Returns a list of::

            {
                "slide_number": int,
                "heading":      str,
                "main_content": str,
                "image_text":   str,
            }
        """
        slides: list[dict] = []
        prs = Presentation(file_path)
        logger.info(
            "Opened PPTX: %d slides, file=%s",
            len(prs.slides),
            file_path,
        )

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_data = PPTExtractor._extract_slide(
                slide, slide_num, file_path
            )
            slides.append(slide_data)

        return slides

    # ── per-slide extraction ────────────────────────────────

    @staticmethod
    def _extract_slide(
        slide: Any, slide_num: int, file_path: str
    ) -> dict:
        heading = ""
        content_parts: list[str] = []
        image_texts: list[str] = []
        img_count = 0
        has_non_picture_shapes = False

        # Identify the title placeholder (if any)
        title_shape_id = None
        shapes = slide.shapes
        if hasattr(shapes, "title") and shapes.title is not None:
            title_shape_id = shapes.title.shape_id

        for shape in shapes:
            shape_type = getattr(shape, "shape_type", None)
            logger.info(
                "Slide %d: shape name=%s, type=%s (%s)",
                slide_num,
                getattr(shape, "name", "?"),
                shape_type,
                type(shape_type).__name__,
            )

            # ── Text frames ─────────────────────────────
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text
                    text = text.strip() if isinstance(text, str) else ""
                    if not text:
                        continue
                    if (
                        title_shape_id is not None
                        and shape.shape_id == title_shape_id
                    ):
                        heading = text
                    else:
                        content_parts.append(text)

            # ── Tables ──────────────────────────────────
            if hasattr(shape, "has_table") and shape.has_table:
                if hasattr(shape, "table") and shape.table:
                    table_text = PPTExtractor._extract_table(
                        shape.table
                    )
                    if table_text:
                        content_parts.append(table_text)

            # ── Images / pictures ───────────────────────
            if (
                PPTExtractor._is_picture_shape(shape)
                and img_count < MAX_IMAGES_PER_SLIDE
            ):
                img_count += 1
                combined = PPTExtractor._process_picture(
                    shape, slide_num
                )
                if combined:
                    image_texts.append(combined)

            # ── Group shapes (may contain nested images) ─
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                has_non_picture_shapes = True
                if hasattr(shape, "shapes"):
                    for child in shape.shapes:
                        if (
                            PPTExtractor._is_picture_shape(child)
                            and img_count < MAX_IMAGES_PER_SLIDE
                        ):
                            img_count += 1
                            combined = PPTExtractor._process_picture(
                                child, slide_num
                            )
                            if combined:
                                image_texts.append(combined)

            # Track non-picture visual shapes (SmartArt, etc.)
            elif (
                shape_type is not None
                and not getattr(shape, "has_text_frame", False)
                and not getattr(shape, "has_table", False)
            ):
                try:
                    st_int = int(shape_type)
                    if st_int not in (0, 1, 6, 17):
                        has_non_picture_shapes = True
                except (ValueError, TypeError):
                    pass

        # Fallback heading
        if not heading:
            if hasattr(shapes, "title") and shapes.title:
                title_text = shapes.title.text
                heading = (
                    title_text.strip()
                    if isinstance(title_text, str)
                    else ""
                )
        if not heading:
            heading = f"Slide {slide_num}"

        # Speaker notes
        if (
            hasattr(slide, "has_notes_slide")
            and slide.has_notes_slide
            and hasattr(slide, "notes_slide")
            and slide.notes_slide
        ):
            ntf = getattr(slide.notes_slide, "notes_text_frame", None)
            if ntf is not None:
                notes_text = getattr(ntf, "text", "")
                notes = (
                    notes_text.strip()
                    if isinstance(notes_text, str)
                    else ""
                )
                if notes:
                    content_parts.append(f"[Notes: {notes}]")

        # Whole-slide caption fallback for SmartArt / vector shapes
        if not image_texts and has_non_picture_shapes:
            logger.info(
                "Slide %d: no picture blobs but has visual shapes "
                "— attempting whole-slide caption fallback",
                slide_num,
            )
            fallback = PPTExtractor._caption_whole_slide(
                file_path, slide_num
            )
            if fallback:
                image_texts.append(f"[Image Description] {fallback}")
        elif not image_texts:
            logger.info(
                "Slide %d: no images or visual shapes found", slide_num
            )

        slide_data = {
            "slide_number": slide_num,
            "heading": clean_text(heading[:200]),
            "main_content": clean_text("\n".join(content_parts)),
            "image_text": clean_text("\n".join(image_texts))[:500],
        }
        logger.info(
            "Slide %d extracted: heading=%s, content_len=%d, "
            "image_text_len=%d, images_found=%d",
            slide_num,
            slide_data["heading"][:50],
            len(slide_data["main_content"]),
            len(slide_data["image_text"]),
            img_count,
        )
        return slide_data

    # ── private helpers ─────────────────────────────────────

    @staticmethod
    def _is_picture_shape(shape: Any) -> bool:
        """
        Detect picture/image shapes.

        Checks for shape type 13 (python-pptx Picture) and also verifies
        the shape actually has an image blob.
        """
        shape_type = getattr(shape, "shape_type", None)
        if shape_type is None:
            return False
        try:
            is_picture = (
                int(shape_type) == 13
                or shape_type == MSO_SHAPE_TYPE.PICTURE
            )
        except Exception:
            is_picture = shape_type == MSO_SHAPE_TYPE.PICTURE

        if is_picture:
            try:
                blob = shape.image.blob
                if blob and len(blob) > 0:
                    return True
                return False
            except Exception:
                return False

        return False

    @staticmethod
    def _process_picture(shape: Any, slide_num: int) -> str:
        """OCR + caption a single picture shape."""
        logger.info(
            "Slide %d: processing picture shape '%s'",
            slide_num,
            getattr(shape, "name", "?"),
        )
        ocr_text = PPTExtractor._ocr_shape_image(shape)
        caption = PPTExtractor._caption_shape_image(shape)
        return PPTExtractor._combine_ocr_and_caption(ocr_text, caption)

    @staticmethod
    def _extract_table(table: Any) -> str:
        """Render a PowerPoint table as pipe-delimited text."""
        rows: list[str] = []
        if not hasattr(table, "rows"):
            return ""
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = getattr(cell, "text", "")
                cells.append(
                    cell_text.strip() if isinstance(cell_text, str) else ""
                )
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    @staticmethod
    def _ocr_shape_image(shape: Any) -> str:
        """Extract the image blob from a Picture shape and run OCR."""
        try:
            image_blob: bytes = shape.image.blob
            if not image_blob:
                logger.info("Shape image blob is empty — skipping OCR")
                return ""
            logger.info(
                "Running OCR on shape image (%d bytes)", len(image_blob)
            )
            text = ocr_image_bytes(image_blob)
            if text:
                logger.info(
                    "OCR extracted %d chars from shape image", len(text)
                )
            return text
        except Exception:
            logger.warning("PPT image OCR failed", exc_info=True)
            return ""

    @staticmethod
    def _caption_shape_image(shape: Any) -> str:
        """Generate an AI caption for a Picture shape using BLIP."""
        temp_path: str | None = None
        try:
            image_blob: bytes = shape.image.blob
            if not image_blob:
                logger.info(
                    "Shape image blob is empty — skipping caption"
                )
                return ""

            logger.info(
                "Extracting image blob (%d bytes) for BLIP captioning",
                len(image_blob),
            )

            with tempfile.NamedTemporaryFile(
                delete=False, suffix=".png"
            ) as tmp:
                tmp.write(image_blob)
                temp_path = tmp.name

            logger.info("Temp image saved to %s — calling BLIP", temp_path)
            caption = generate_image_caption(temp_path)
            if caption:
                logger.info(
                    "Generated BLIP caption: %s", caption[:80]
                )
            else:
                logger.warning(
                    "BLIP returned empty caption (%d bytes)",
                    len(image_blob),
                )
            return caption
        except Exception:
            logger.warning("PPT image captioning failed", exc_info=True)
            return ""
        finally:
            if temp_path and os.path.exists(temp_path):
                try:
                    os.unlink(temp_path)
                except OSError:
                    logger.debug(
                        "Failed to delete temp image: %s", temp_path
                    )

    @staticmethod
    def _combine_ocr_and_caption(ocr_text: str, caption: str) -> str:
        """Merge OCR text and BLIP caption into a single string."""
        parts: list[str] = []
        if ocr_text:
            parts.append(f"[OCR Text] {ocr_text}")
        if caption:
            parts.append(f"[Image Description] {caption}")
        return "\n".join(parts)

    @staticmethod
    def _caption_whole_slide(
        file_path: str, slide_number: int
    ) -> str:
        """
        Render a full slide to an image and caption it with BLIP.

        Fallback for decks where visuals are vector/SmartArt/grouped
        shapes without extractable image blobs.
        """
        temp_dir = tempfile.mkdtemp(prefix="slide_reader_ppt_fallback_")

        try:
            input_path = Path(file_path)
            soffice_cmd = [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                temp_dir,
                str(input_path),
            ]
            completed = subprocess.run(
                soffice_cmd,
                check=False,
                capture_output=True,
                text=True,
            )

            if completed.returncode != 0:
                logger.debug(
                    "Slide fallback conversion failed (slide=%d): %s",
                    slide_number,
                    (
                        completed.stderr.strip()
                        if completed.stderr
                        else "unknown error"
                    ),
                )
                return ""

            pdf_path = os.path.join(
                temp_dir, f"{input_path.stem}.pdf"
            )
            if not os.path.exists(pdf_path):
                logger.debug(
                    "Fallback PDF not found for slide %d", slide_number
                )
                return ""

            with fitz.open(pdf_path) as doc:
                page_idx = slide_number - 1
                if page_idx < 0 or page_idx >= len(doc):
                    return ""

                page = doc[page_idx]
                pix = page.get_pixmap(
                    matrix=fitz.Matrix(2.0, 2.0)
                )
                png_path = os.path.join(
                    temp_dir, f"slide_{slide_number}.png"
                )
                pix.save(png_path)

            return generate_image_caption(png_path)

        except FileNotFoundError:
            logger.debug(
                "`soffice` not available; "
                "skipping whole-slide caption fallback"
            )
            return ""
        except Exception:
            logger.debug(
                "Whole-slide caption fallback failed", exc_info=True
            )
            return ""
        finally:
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except Exception:
                logger.debug(
                    "Failed to remove fallback temp dir: %s", temp_dir
                )
